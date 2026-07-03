from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

import os
from datetime import datetime

from slide_models import SlideModel, BoxModel, PresentationModel


class PptService:
    DEFAULT_COLORS = {
        "HEAD": RGBColor(0, 15, 45),
        "ASSISTANT": RGBColor(26, 153, 171),
        "UNIT": RGBColor(189, 215, 255),
        "DOTTED": RGBColor(89, 89, 89),
        "HIGHLIGHT": RGBColor(0, 204, 102),
        "LINE": RGBColor(90, 90, 90)
    }

    def __init__(self, template_dir: str, output_dir: str, base_url: str):
        self.template_dir = template_dir
        self.output_dir = output_dir
        self.base_url = base_url
        self.COLORS = dict(self.DEFAULT_COLORS)

    @staticmethod
    def _parse_hex_color(value):
        if not value or not isinstance(value, str):
            return None
        hex_str = value.strip().lstrip("#")
        if len(hex_str) != 6:
            return None
        try:
            return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
        except ValueError:
            return None

    def _apply_color_overrides(self, colors):
        self.COLORS = dict(self.DEFAULT_COLORS)
        if not colors:
            return
        for key, value in colors.items():
            key = str(key).upper()
            if key not in self.DEFAULT_COLORS:
                continue
            parsed = self._parse_hex_color(value)
            if parsed is not None:
                self.COLORS[key] = parsed

    def generate_presentation(self, model, template_name, output_name=None, colors=None):
        self._apply_color_overrides(colors)
        template_path = os.path.join(self.template_dir, f"{template_name}.pptx")
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template not found: {template_path}")
        prs = Presentation(template_path)
        for i in range(len(prs.slides) - 1, -1, -1):
            r_id = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(r_id)
            del prs.slides._sldIdLst[i]
        org_id_to_slide = {}
        pending_links = []

        pending_back_buttons = []
        primary_slide = None

        for slide_model in model.slides:
            if getattr(slide_model, "type", None) == "org_chart":
                is_primary = primary_slide is None
                slide, link_shapes = self._add_org_chart_slide(
                    prs, slide_model, is_primary=is_primary
                )

                if is_primary:
                    primary_slide = slide

                slide_org_id = getattr(slide_model, "org_id", None)
                if slide_org_id is not None:
                    slide_org_id = str(slide_org_id)
                    if slide_org_id not in org_id_to_slide:
                        org_id_to_slide[slide_org_id] = slide
                pending_links.extend(link_shapes)

                if not is_primary:
                    back_shape = self._add_back_button(slide)
                    pending_back_buttons.append(back_shape)

        for shape, target_org_id in pending_links:
            target_slide = org_id_to_slide.get(str(target_org_id))
            if target_slide is not None:
                self._set_internal_hyperlink(shape, target_slide)

        # Wire the back buttons to the primary slide.
        if primary_slide is not None:
            for back_shape in pending_back_buttons:
                self._set_internal_hyperlink(back_shape, primary_slide)
        filename = self._generate_filename(output_name)
        output_path = os.path.join(self.output_dir, filename)
        os.makedirs(self.output_dir, exist_ok=True)
        prs.save(output_path)
        return self._build_file_url(filename)

    def _add_org_chart_slide(self, prs, slide_model, is_primary=False):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                sp = shape._sp
                sp.getparent().remove(sp)
        title_shape = slide.shapes.add_textbox(Cm(1.3), Cm(0.76), Cm(22.86), Cm(2.54))
        tf = title_shape.text_frame
        tf.text = f"{slide_model.title} {slide_model.hierarchy_level}".strip()
        title_paragraph = tf.paragraphs[0]
        title_paragraph.font.name = "Nokia Pure Headline Light"
        title_paragraph.font.size = Pt(28)
        title_paragraph.font.color.rgb = RGBColor(0, 102, 255)

        col_width = Cm(8)
        start_y = Cm(3.05)
        row_spacing = Cm(1.59)
        box_width = Cm(6.93)
        box_height = Cm(1.33)
        manager_x = Cm(7.62)
        manager_y = Cm(15.75)
        manager_width = Cm(6.93)
        manager_height = Cm(1.33)

        all_box_positions = []
        for i, column in enumerate(slide_model.columns):
            current_x = Cm(3.8) + (i * (col_width + Cm(0.76)))
            column_positions = []
            for j, unit in enumerate(column):
                unit_y = start_y + (j * row_spacing)
                column_positions.append({"box": unit, "x": current_x, "y": unit_y,
                                         "width": box_width, "height": box_height})
            if column_positions:
                all_box_positions.append(column_positions)

        if all_box_positions:
            self._draw_connection_lines(slide=slide, all_box_positions=all_box_positions,
                manager_x=manager_x, manager_y=manager_y, manager_width=manager_width,
                box_height=box_height, start_y=start_y)

        link_shapes = []
        for column_positions in all_box_positions:
            for box_position in column_positions:
                shape = self._draw_box(slide, box_position["x"], box_position["y"],
                    box_position["box"], width=box_position["width"], height=box_position["height"])
                target_org_id = getattr(box_position["box"], "target_org_id", None)
                if target_org_id is not None:
                    link_shapes.append((shape, str(target_org_id)))

        self._draw_box(slide, manager_x, manager_y, slide_model.head_of_unit,
                       width=manager_width, height=manager_height)

        if slide_model.assistant:
            self._draw_box(slide, Cm(14.23), Cm(15.75), slide_model.assistant,
                           width=Cm(6.93), height=Cm(1.33))

        self._add_legend(slide)
        self._add_footer(slide)
        return slide, link_shapes

    def _draw_connection_lines(self, slide, all_box_positions, manager_x, manager_y,
                               manager_width, box_height, start_y):
        manager_center_x = int(manager_x + manager_width / 2)
        manager_center_y = int(manager_y + Cm(1.33) / 2)
        main_line_y = start_y + box_height / 2
        if not all_box_positions:
            return
        column_spine_x_values = []
        for column_positions in all_box_positions:
            first_box = column_positions[0]
            spine_x = int(first_box["x"] - Cm(0.8))
            column_spine_x_values.append(spine_x)
        leftmost_x = min(column_spine_x_values)
        rightmost_x = max(column_spine_x_values)
        bottom_connector_y = manager_y - box_height / 2
        self._draw_horizontal_line(slide, leftmost_x, manager_center_x, bottom_connector_y)
        self._draw_horizontal_line(slide, manager_center_x, rightmost_x, bottom_connector_y)
        self._draw_vertical_line(slide, manager_center_x, bottom_connector_y, manager_y)
        for col_idx, column_positions in enumerate(all_box_positions):
            spine_x = column_spine_x_values[col_idx]
            self._draw_vertical_line(slide, spine_x, main_line_y, bottom_connector_y)
            for box_position in column_positions:
                box_center_y = int(box_position["y"] + box_position["height"] / 2)
                box_left_x = int(box_position["x"])
                self._draw_horizontal_line(slide, spine_x, box_left_x, box_center_y)

    def _draw_vertical_line(self, slide, x, y_top, y_bottom):
        x = int(x); y_top = int(y_top); y_bottom = int(y_bottom)
        if y_top > y_bottom:
            y_top, y_bottom = y_bottom, y_top
        line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x, y_top, x, y_bottom)
        line.line.fill.solid()
        line.line.fill.fore_color.rgb = self.COLORS["LINE"]
        line.line.color.rgb = self.COLORS["LINE"]
        line.line.width = Pt(1.2)
        return line

    def _draw_horizontal_line(self, slide, x_left, x_right, y):
        x_left = int(x_left); x_right = int(x_right); y = int(y)
        if x_left > x_right:
            x_left, x_right = x_right, x_left
        line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x_left, y, x_right, y)
        line.line.fill.solid()
        line.line.fill.fore_color.rgb = self.COLORS["LINE"]
        line.line.color.rgb = self.COLORS["LINE"]
        line.line.width = Pt(1.2)
        return line

    def _draw_line(self, slide, x1, y1, x2, y2):
        line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
            int(x1), int(y1), int(x2), int(y2))
        line.line.fill.solid()
        line.line.fill.fore_color.rgb = self.COLORS["LINE"]
        line.line.color.rgb = self.COLORS["LINE"]
        line.line.width = Pt(1.2)
        return line

    def _add_legend(self, slide):
        """
        Small color legend in the bottom-right corner, near the Nokia logo,
        so it does not overlap the chart boxes. Swatch colors follow the
        active palette (presets / custom colors are reflected here too).
        """
        entries = [
            ("Head of unit", "HEAD", False),
            ("Unit", "UNIT", False),
            ("Connectors", "LINE", True),
        ]
        # Bottom-right corner, above the Nokia logo (slide ~33.87 cm wide).
        legend_x = Cm(29.0)
        legend_y = Cm(13.7)
        row_height = Cm(0.62)
        swatch_size = Cm(0.42)
        label_x = legend_x + swatch_size + Cm(0.2)
        label_width = Cm(4.0)

        for i, (label, color_key, is_line) in enumerate(entries):
            row_y = legend_y + (i * row_height)
            if is_line:
                line_y = int(row_y + swatch_size / 2)
                line = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                    int(legend_x), line_y, int(legend_x + swatch_size), line_y)
                line.line.color.rgb = self.COLORS[color_key]
                line.line.width = Pt(2)
            else:
                swatch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                    legend_x, row_y, swatch_size, swatch_size)
                swatch.fill.solid()
                swatch.fill.fore_color.rgb = self.COLORS[color_key]
                swatch.line.fill.background()
            label_box = slide.shapes.add_textbox(label_x, row_y - Cm(0.1), label_width, Cm(0.6))
            ltf = label_box.text_frame
            ltf.word_wrap = False
            ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
            lp = ltf.paragraphs[0]
            lp.text = label
            lp.font.name = "Nokia Pure Text Light"
            lp.font.size = Pt(9)
            lp.font.color.rgb = RGBColor(64, 64, 64)

    def _add_back_button(self, slide):
        """
        Draw a "Back to primary organization" button in the bottom-left corner,
        near the page number. Returns the shape so the caller can attach the
        internal hyperlink to the primary slide.
        """
        left = Cm(1.0)
        top = Cm(16.0)
        width = Cm(5.6)
        height = Cm(0.85)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 102, 255)
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Cm(0.2)
        tf.margin_right = Cm(0.2)
        tf.margin_top = Cm(0.05)
        tf.margin_bottom = Cm(0.05)

        p = tf.paragraphs[0]
        p.text = "\u2190 Back to primary organization"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Nokia Pure Text Light"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        return shape

    def _add_footer(self, slide):
        left = Cm(5.57); top = Cm(17.24); width = Cm(22.80); height = Cm(1.54)
        footer_shape = slide.shapes.add_textbox(left, top, width, height)
        tf = footer_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ("This document aims to support the business groups and corporate functions "
            "in Nokia subject to necessary legal procedures and approvals. Nothing herein "
            "should be deemed to indicate that any final decision has been made which "
            "otherwise may require information and/or consultation with relevant employee "
            "representative body/ies, where applicable.")
        p.font.size = Pt(10.7)
        p.font.name = "Nokia Pure Text Light"
        p.font.color.rgb = RGBColor(128, 128, 128)
        p.alignment = PP_ALIGN.LEFT

    def _draw_box(self, slide, x, y, box_data, width=Cm(6.93), height=Cm(1.33)):
        CHARS_PER_LINE = 40
        title_len = len(box_data.title or "")
        title_lines = max(1, -(-title_len // CHARS_PER_LINE))
        sub_len = len(box_data.subtitle or "")
        sub_lines = max(1, -(-sub_len // CHARS_PER_LINE))
        total_lines = title_lines + sub_lines
        if total_lines > 2:
            calculated_height = Cm(0.4) + Cm(0.45 * total_lines)
            height = max(height, calculated_height)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, height)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        if box_data.is_changed:
            shape.line.color.rgb = self.COLORS["HIGHLIGHT"]
            shape.line.width = Pt(2.5)
        else:
            shape.line.fill.background()
        shape.fill.solid()
        if box_data.box_type == "manager":
            shape.fill.fore_color.rgb = self.COLORS["HEAD"]
        elif box_data.box_type == "assistant":
            shape.fill.fore_color.rgb = self.COLORS["ASSISTANT"]
        elif box_data.box_type == "dotted":
            shape.fill.fore_color.rgb = self.COLORS["DOTTED"]
        else:
            shape.fill.fore_color.rgb = self.COLORS["UNIT"]
        tf.margin_left = Cm(0.25)
        tf.margin_right = Cm(0.15)
        tf.margin_top = Cm(0.05)
        tf.margin_bottom = Cm(0.05)
        p1 = tf.paragraphs[0]
        p1.text = box_data.title
        p1.font.name = "Nokia Pure Headline"
        p1.alignment = PP_ALIGN.LEFT
        p1.font.bold = True
        p1.font.size = Pt(12)
        if box_data.box_type in ["manager", "assistant", "dotted"]:
            p1.font.color.rgb = RGBColor(255, 255, 255)
        else:
            p1.font.color.rgb = RGBColor(0, 0, 0)
        p2 = tf.add_paragraph()
        p2.text = box_data.subtitle
        p2.font.name = "Nokia Pure Headline"
        p2.font.size = Pt(12)
        p2.font.color.rgb = p1.font.color.rgb
        return shape

    def _set_internal_hyperlink(self, shape, target_slide):
        slide_part = shape.part
        rId = slide_part.relate_to(target_slide.part,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide")
        cNvPr = shape._element._nvXxPr.cNvPr
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        existing = cNvPr.find(f"{{{ns}}}hlinkClick")
        if existing is not None:
            cNvPr.remove(existing)
        hlink = cNvPr.makeelement(f"{{{ns}}}hlinkClick", {})
        hlink.set(f"{{{r_ns}}}id", rId)
        hlink.set("action", "ppaction://hlinksldjump")
        cNvPr.append(hlink)

    def _generate_filename(self, output_name):
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        if output_name:
            clean_name = output_name.replace(".pptx", "")
            return f"{clean_name}_{timestamp}.pptx"
        return f"presentation_{timestamp}.pptx"

    def _build_file_url(self, filename):
        return f"{self.base_url}/{filename}"